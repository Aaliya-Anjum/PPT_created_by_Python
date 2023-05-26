from pptx import Presentation

def make_ppt(file_path, file_path2, ppt_path):
    prs = Presentation()

    with open(file_path, 'r') as sample1:
        content = sample1.read()

    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = 'Slide 1'  
    content_placeholder = slide.placeholders[1]
    content_placeholder.text = content

    with open(file_path2, 'r') as sample2:
        content1 = sample2.read()

    slide_layout1 = prs.slide_layouts[1] 
    slide1 = prs.slides.add_slide(slide_layout1)
    title1 = slide1.shapes.title
    title1.text = 'Slide 2' 
    content_placeholder1 = slide1.placeholders[1]  
    content_placeholder1.text = content1

   # prs.save(ppt_path)
    prs.save('presentation.pptx')

    print('PPT file created successfully.')

file_path = 'D:\Aaliya\Presentation\sample.txt'
file_path2 = 'D:\Aaliya\Presentation\sample2.txt'
ppt_path = r'D:\Aaliya\Presentation\exa.pptx'

make_ppt(file_path, file_path2, ppt_path)
